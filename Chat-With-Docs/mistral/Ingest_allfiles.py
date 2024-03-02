import os
from pptx import Presentation
from docx import Document as DocxDocument
import box
import yaml
from langchain.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.document_loaders import PyPDFLoader, DirectoryLoader, TextLoader
from langchain.embeddings import HuggingFaceEmbeddings
import re

def convert_latex_to_text(latex_content):
    latex_content = re.sub(r'%.*?\n', '', latex_content)
    plain_text = re.sub(r'\\[^{}]+{|}|[\[\]\{\}]', '', latex_content)
    return plain_text

def read_latex_file(input_file):
    with open(input_file, 'r', encoding='utf-8') as f:
        latex_content = f.read()
    return latex_content

def write_plain_text_file(output_file, plain_text):
    with open(output_file, 'w', encoding='utf-8') as f:
        f.write(plain_text)

with open('mistral/config.yml', 'r', encoding='utf8') as ymlfile:
    cfg = box.Box(yaml.safe_load(ymlfile))

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
    return text

def extract_text_from_docx(path):
    doc = DocxDocument(path)
    return ' '.join([paragraph.text for paragraph in doc.paragraphs])

def run_ingest():
    pdf_loader = DirectoryLoader(cfg.DATA_PATH,
                                 glob='*.pdf',
                                 loader_cls=PyPDFLoader)
    
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=cfg.CHUNK_SIZE,
                                                   chunk_overlap=cfg.CHUNK_OVERLAP,length_function = len,
    add_start_index = True,)

    documents = pdf_loader.load()
    docs=[]
    ctr=0
    for i in documents:
      doc=i.page_content
      docs+=text_splitter.create_documents([doc])

    txt_docs=[]
    for file in os.listdir(cfg.DATA_PATH):
        if file.endswith(".pptx"):
            text = extract_text_from_pptx(os.path.join(cfg.DATA_PATH, file))
            with open(f'{file}.txt', 'w') as f:
                f.write(text)
            txt_loader = TextLoader(f'{file}.txt')
            txt_document = txt_loader.load()
            for i in txt_document:
              doc=i.page_content
              txt_docs+=text_splitter.create_documents([doc])

    docx_docs=[]
    for file in os.listdir(cfg.DATA_PATH):
        if file.endswith(".docx"):
            print('hi')
            text = extract_text_from_docx(os.path.join(cfg.DATA_PATH, file))
            with open(f'{file}.txt', 'w') as f:
                f.write(text)
            txt_loader = TextLoader(f'{file}.txt')
            docx_document = txt_loader.load()
            for i in docx_document:
              doc=i.page_content
              docx_docs+=text_splitter.create_documents([doc])

    text_docs=[]
    for file in os.listdir(cfg.DATA_PATH):
        if file.endswith(".txt"):
            txt_loader = TextLoader(f'{file}.txt')
            txt_document = txt_loader.load()
            for i in txt_document:
                doc = i.page_content
                text_docs+=text_splitter.create_documents([doc])

    texts=[doc.page_content for doc in docs] + [doc.page_content for doc in txt_docs] + [doc.page_content for doc in docx_docs] + [doc.page_content for doc in text_docs]
    embeddings = HuggingFaceEmbeddings(model_name=cfg.EMBEDDINGS,
                                       model_kwargs={'device': 'cpu'})
    metadata=[{"source":text} for text in texts]
    vectorstore = FAISS.from_texts(texts, embeddings,metadatas=metadata)
    vectorstore.save_local(cfg.DB_FAISS_PATH)

if __name__ == "__main__":
    run_ingest()